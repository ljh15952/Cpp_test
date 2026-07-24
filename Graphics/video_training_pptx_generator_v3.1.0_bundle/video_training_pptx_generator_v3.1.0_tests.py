from __future__ import annotations

import asyncio
import ast
import importlib.util
import json
import sys
import tempfile
import types
from pathlib import Path
from typing import Any

from PIL import Image, ImageDraw
from pptx import Presentation

TOOL_PATH = Path('/mnt/data/video_training_pptx_generator_v3.1.0.py')
PROMPT_PATH = Path('/mnt/data/video_training_pptx_generator_prompt_v3.1.0.md')


def load_module():
    spec = importlib.util.spec_from_file_location('video_training_pptx_generator_v31', TOOL_PATH)
    if spec is None or spec.loader is None:
        raise RuntimeError('module spec failure')
    module = importlib.util.module_from_spec(spec)
    sys.modules[spec.name] = module
    spec.loader.exec_module(module)
    return module


def assert_true(condition: bool, message: str) -> None:
    if not condition:
        raise AssertionError(message)


def sample_segments() -> list[dict[str, Any]]:
    return [
        {"index": 0, "start": 0.0, "end": 8.0, "text": "この講座では売上データの集計方法を説明します。"},
        {"index": 1, "start": 8.0, "end": 22.0, "text": "挿入タブからピボットテーブルを選択します。"},
        {"index": 2, "start": 22.0, "end": 38.0, "text": "年月と商品分類を配置すると集計結果が表示されます。"},
        {"index": 3, "start": 38.0, "end": 50.0, "text": "最後に結果を確認して講座をまとめます。"},
    ]


def sample_plan() -> dict[str, Any]:
    notes = {
        "explanation_points": ["ピボットテーブルの目的を説明します。"],
        "demo_steps": ["挿入タブから操作を実演します。"],
        "questions": ["集計結果を確認できましたか。"],
        "cautions": ["元データの見出しを確認します。"],
        "estimated_minutes": 1.5,
        "evidence": ["音声セグメント1～2"],
        "review_items": [],
    }
    return {
        "deck_title": "Excel集計講座",
        "subtitle": "ピボットテーブルの基本操作",
        "target_audience": "初学者",
        "course_objectives": ["ピボットテーブルで集計できる"],
        "slides": [
            {
                "slide_number": 1,
                "slide_id": "create-pivot",
                "title": "ピボットテーブルを作成する",
                "type": "operation",
                "objective": "元データから集計表を作成する",
                "bullets": [],
                "steps": [
                    {"step_no": 1, "target": "挿入タブ", "action": "選択する", "value": "ピボットテーブル", "expected_result": "作成画面が表示される"},
                    {"step_no": 2, "target": "フィールド", "action": "配置する", "value": "年月・商品分類", "expected_result": "集計結果が表示される"},
                ],
                "source_segment_indexes": [1, 2],
                "start_seconds": 8.0,
                "end_seconds": 38.0,
                "frame": {"timestamp_seconds": 24.0, "phase": "during", "purpose": "操作画面"},
                "speaker_notes": notes,
            },
            {
                "slide_number": 2,
                "slide_id": "confirm-result",
                "title": "集計結果を確認する",
                "type": "summary",
                "objective": "作成した集計表の結果を確認する",
                "bullets": ["年月と商品分類ごとの集計結果を確認します。", "必要に応じて表示項目を調整します。"],
                "steps": [],
                "source_segment_indexes": [2, 3],
                "start_seconds": 22.0,
                "end_seconds": 50.0,
                "frame": {"timestamp_seconds": 42.0, "phase": "after", "purpose": "結果画面"},
                "speaker_notes": notes,
            },
        ],
    }


async def run() -> list[str]:
    module = load_module()
    tool = module.Tools()
    passed: list[str] = []

    # 1. Syntax and public surface.
    tree = ast.parse(TOOL_PATH.read_text(encoding='utf-8'))
    tools_cls = next(n for n in tree.body if isinstance(n, ast.ClassDef) and n.name == 'Tools')
    public = sorted(
        n.name for n in tools_cls.body
        if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef)) and not n.name.startswith('_')
    )
    assert_true(public == ['debug_ping', 'generate_training_material'], f'public actions: {public}')
    passed.append('public actions limited to debug_ping and generate_training_material')

    # 2. Evaluation/release/repair pipeline remains removed.
    source = TOOL_PATH.read_text(encoding='utf-8')
    forbidden = ['quality_score', 'release_status', 'repair_training_material', '_semantic_validate_plan', '_validate_pptx', 'major_error_count', 'validation.json']
    found = [token for token in forbidden if token in source]
    assert_true(not found, f'forbidden evaluation tokens: {found}')
    passed.append('evaluation, release blocking, and targeted repair remain removed')

    # 3. All self-method references resolve.
    methods = {n.name for n in tools_cls.body if isinstance(n, (ast.FunctionDef, ast.AsyncFunctionDef))}
    refs = set()
    for node in ast.walk(tools_cls):
        if isinstance(node, ast.Attribute) and isinstance(node.value, ast.Name):
            if node.value.id in {'self', 'cls', 'Tools'} and node.attr.startswith('_'):
                refs.add(node.attr)
    missing = sorted(refs - methods - {'_json_retry_lock', '_generation_guard_lock', '_generation_inflight', '_generation_result_cache'})
    assert_true(not missing, f'missing method refs: {missing}')
    passed.append('all internal method references resolve')

    # 4. debug_ping contract.
    ping = json.loads(await tool.debug_ping())
    assert_true(ping['version'] == '3.1.0', 'wrong version')
    assert_true(ping['available_actions'] == ['debug_ping', 'generate_training_material'], 'wrong actions')
    assert_true(ping['settings']['editorial_mode'] == 'batched_with_local_fallback', 'wrong editor mode')
    assert_true(ping['settings']['continue_on_editor_error'] is True, 'editor fallback disabled')
    passed.append('debug_ping v3.1 resilient editorial contract')

    # 5. Top-level JSON arrays are accepted and mapped to the expected list key.
    parsed = tool._parse_json_payload('[{"title":"A"}]')
    assert_true(isinstance(parsed, list) and parsed[0]['title'] == 'A', 'array parse failure')
    saved_modules = {name: sys.modules.get(name) for name in (
        'open_webui', 'open_webui.models', 'open_webui.models.users',
        'open_webui.utils', 'open_webui.utils.chat',
    )}
    open_webui = types.ModuleType('open_webui')
    models_mod = types.ModuleType('open_webui.models')
    users_mod = types.ModuleType('open_webui.models.users')
    utils_mod = types.ModuleType('open_webui.utils')
    chat_mod = types.ModuleType('open_webui.utils.chat')
    class FakeUsers:
        @staticmethod
        def get_user_by_id(user_id):
            return {'id': user_id}
    async def fake_completion(request, form_data, user_model, **kwargs):
        return {'choices': [{'message': {'content': '[{"title":"A"}]'}}]}
    users_mod.Users = FakeUsers
    chat_mod.generate_chat_completion = fake_completion
    sys.modules.update({
        'open_webui': open_webui,
        'open_webui.models': models_mod,
        'open_webui.models.users': users_mod,
        'open_webui.utils': utils_mod,
        'open_webui.utils.chat': chat_mod,
    })
    try:
        mapped = await tool._call_internal_llm_json(
            system='system', prompt='prompt', request=object(), user={'id': 'u'},
            model={'id': 'm'}, metadata={}, list_key='slides',
        )
    finally:
        for name, previous in saved_modules.items():
            if previous is None:
                sys.modules.pop(name, None)
            else:
                sys.modules[name] = previous
    assert_true(mapped == {'slides': [{'title': 'A'}]}, str(mapped))
    passed.append('top-level JSON array compatibility and list-key mapping')

    # 6. Planner format failures use rule-based fallback even when LLM is required.
    tool = module.Tools()
    async def broken_planner(**kwargs):
        raise RuntimeError('内部LLMのJSON応答を期待する構造へ変換できません。')
    tool._call_internal_llm_json = broken_planner  # type: ignore[method-assign]
    chapters = await tool._build_chapters(
        segments=sample_segments(), title='Excel集計講座', audience='初学者', language='日本語',
        request=object(), user={'id': 'u'}, model={'id': 'm'}, metadata={}, emitter=None,
    )
    assert_true(bool(chapters), 'chapter fallback failed')
    fallback_plan = await tool._build_slide_plan(
        chapters=chapters, segments=sample_segments(), title='Excel集計講座', audience='初学者',
        language='日本語', min_slides=1, max_slides=10, request=object(), user={'id': 'u'},
        model={'id': 'm'}, metadata={}, emitter=None,
    )
    assert_true(bool(fallback_plan.get('slides')), 'slide-plan fallback failed')
    passed.append('planner JSON-format failure continues with rule-based plan')

    # 7. Minimum slide target does not force thin slides.
    normalized = tool._normalize_slide_plan(sample_plan(), sample_segments(), 'Excel集計講座', '初学者', 8, 24)
    assert_true(len(normalized['slides']) == 2, 'minimum count forced extra slides')
    passed.append('minimum slide count is a preference, not a forced floor')

    # 8. Editor failure falls back without aborting generation.
    async def broken_llm(**kwargs):
        raise RuntimeError('内部LLMのJSON応答がオブジェクトではありません。')
    tool._call_internal_llm_json = broken_llm  # type: ignore[method-assign]
    polished = await tool._polish_slide_plan(
        plan=sample_plan(), segments=sample_segments(), title='Excel集計講座',
        audience='初学者', language='日本語', preferred_min_slides=6,
        preferred_max_slides=24, request=object(), user={'id': 'u'}, model={'id': 'm'}, metadata={},
    )
    assert_true(len(polished['slides']) == 2, 'fallback lost slides')
    assert_true(polished.get('_editor_fallback_count') == 1, str(polished.get('_editor_fallback_count')))
    passed.append('editor JSON failure uses local fallback and continues')

    # 9. Batch editor preserves IDs/evidence and allows edited text.
    tool = module.Tools()
    async def good_llm(**kwargs):
        slides = sample_plan()['slides']
        edited = json.loads(json.dumps(slides, ensure_ascii=False))
        edited[0]['title'] = 'ピボットテーブル作成手順'
        edited[0]['source_segment_indexes'] = [0]  # must be ignored by safe merge
        return {'slides': edited}
    tool._call_internal_llm_json = good_llm  # type: ignore[method-assign]
    polished = await tool._polish_slide_plan(
        plan=sample_plan(), segments=sample_segments(), title='Excel集計講座',
        audience='初学者', language='日本語', preferred_min_slides=1,
        preferred_max_slides=24, request=object(), user={'id': 'u'}, model={'id': 'm'}, metadata={},
    )
    assert_true(polished['slides'][0]['title'].startswith('ピボットテーブル作成手順'), 'edited title not applied')
    assert_true(polished['slides'][0]['source_segment_indexes'] == [1, 2], 'evidence indexes changed')
    passed.append('batched editor safely merges editable fields only')

    # 10. Actual PPTX creation.
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        img1 = root / 'frame1.png'
        img2 = root / 'frame2.png'
        for idx, path in enumerate((img1, img2), start=1):
            image = Image.new('RGB', (1280, 720), (235, 239, 245))
            draw = ImageDraw.Draw(image)
            draw.rectangle((80, 80, 1200, 640), outline=(80, 95, 115), width=4)
            draw.text((120, 120), f'Sample screen {idx}', fill=(20, 30, 45))
            image.save(path)
        pptx_path = root / 'sample.pptx'
        records = [
            {"slide_number": 1, "image_path": str(img1), "actual_frame_pts_seconds": 24.0},
            {"slide_number": 2, "image_path": str(img2), "actual_frame_pts_seconds": 42.0},
        ]
        tool._create_pptx(pptx_path, sample_plan(), records, 'sample.mp4', '初学者')
        assert_true(pptx_path.is_file() and pptx_path.stat().st_size > 1024, 'pptx not created')
        prs = Presentation(str(pptx_path))
        assert_true(len(prs.slides) == 3, f'unexpected slide count: {len(prs.slides)}')
        passed.append('actual PPTX generation with adaptive layouts')

    # 11. End-to-end generation with heavy operations mocked.
    with tempfile.TemporaryDirectory() as td:
        root = Path(td)
        video = root / 'demo.mp4'
        video.write_bytes(b'video-placeholder')
        tool = module.Tools()
        tool.valves.output_dir = str(root / 'out')
        tool._validate_dependencies = lambda **kwargs: None  # type: ignore[method-assign]
        tool._validate_internal_llm_context = lambda **kwargs: None  # type: ignore[method-assign]
        tool._resolve_video_file = lambda files: (video, 'demo.mp4')  # type: ignore[method-assign]
        tool._probe_media = lambda path: module.MediaInfo(50.0, 1280, 720, 30.0, 0.0, True, 1)  # type: ignore[method-assign]
        def fake_extract_audio(video_path, audio_path, stream_index):
            audio_path.write_bytes(b'0' * 2048)
        tool._extract_audio = fake_extract_audio  # type: ignore[method-assign]
        tool._transcribe_audio = lambda path: (sample_segments(), {"language": "ja"})  # type: ignore[method-assign]
        tool._probe_all_frame_pts = lambda path, media: [0.0, 8.0, 24.0, 42.0, 49.0]  # type: ignore[method-assign]
        async def fake_chapters(**kwargs):
            return [{"title": "集計", "type": "operation", "source_segment_indexes": [0,1,2,3], "start_seconds": 0.0, "end_seconds": 50.0, "summary": "集計操作", "key_points": [], "operations": [], "cautions": [], "uncertainties": []}]
        async def fake_plan(**kwargs): return sample_plan()
        async def fake_polish(**kwargs): return sample_plan()
        async def fake_frames(**kwargs):
            frames_dir = kwargs['frames_dir']
            frames_dir.mkdir(parents=True, exist_ok=True)
            records = []
            for idx, pts in ((1, 24.0), (2, 42.0)):
                path = frames_dir / f'slide_{idx:03d}.png'
                Image.new('RGB', (1280, 720), (240, 242, 246)).save(path)
                records.append({"slide_number": idx, "slide_id": f's{idx}', "image_path": str(path), "actual_frame_pts_seconds": pts})
            return records
        async def fake_attach(**kwargs): return {"id": "file123", "url": "/files/file123", "emitted": True}
        tool._build_chapters = fake_chapters  # type: ignore[method-assign]
        tool._build_slide_plan = fake_plan  # type: ignore[method-assign]
        tool._polish_slide_plan = fake_polish  # type: ignore[method-assign]
        tool._prepare_slide_frames = fake_frames  # type: ignore[method-assign]
        tool._register_and_emit_generated_file = fake_attach  # type: ignore[method-assign]
        result = json.loads(await tool.generate_training_material(
            __files__=[{"path": str(video), "name": "demo.mp4"}],
            __user__={"id": "u"}, __model__={"id": "m"}, __metadata__={}, __request__=object(),
        ))
        assert_true(result['status'] == 'success' and result['attached'] is True, str(result))
        assert_true(result['generation_mode'] == 'resilient_batched_editorial', str(result))
        output = Path(result['output_file_path'])
        assert_true(output.is_file(), 'end-to-end output missing')
        children = list(output.parent.iterdir())
        assert_true(children == [output], f'intermediate files remain: {children}')
        passed.append('mocked end-to-end generation and PPTX-only cleanup')

    # 12. Duplicate calls for the same message reuse the first result.
    tool = module.Tools()
    call_count = 0
    async def fake_impl(**kwargs):
        nonlocal call_count
        call_count += 1
        await asyncio.sleep(0.01)
        return json.dumps({"status": "error", "message": "first result"}, ensure_ascii=False)
    tool._generate_training_material_impl = fake_impl  # type: ignore[method-assign]
    metadata = {"chat_id": "chat-1", "message_id": "msg-1"}
    first, second = await asyncio.gather(
        tool.generate_training_material(__files__=[{"id": "f1", "name": "demo.mp4"}], __metadata__=metadata),
        tool.generate_training_material(__files__=[{"id": "f1", "name": "demo.mp4"}], __metadata__=metadata),
    )
    assert_true(call_count == 1, f'duplicate execution count: {call_count}')
    payloads = [json.loads(first), json.loads(second)]
    assert_true(any(p.get('duplicate_request_reused') for p in payloads), str(payloads))
    passed.append('duplicate calls reuse the first generation result')

    # 13. Error result explicitly prohibits automatic retry.
    assert_true('"automatic_retry_allowed": False' in source, 'error result retry guard missing')
    prompt = PROMPT_PATH.read_text(encoding='utf-8')
    assert_true('同じ利用者メッセージの処理中に`generate_training_material`を再度呼び出してはいけません' in prompt, 'prompt retry prohibition missing')
    assert_true('エラー結果を表示した時点で、そのターンを直ちに終了' in prompt, 'prompt stop rule missing')
    passed.append('automatic duplicate generation is prohibited')

    # 14. Prompt contains resilient JSON/editor fallback contract.
    assert_true('トップレベルがJSON配列' in prompt, 'array compatibility missing')
    assert_true('ローカル編集フォールバック' in prompt, 'editor fallback missing')
    assert_true('汎用説明動画PPTX自動生成_v3.1.0' in prompt, 'tool version missing')
    passed.append('system prompt v3.1 resilient routing and fallback contract')

    return passed


if __name__ == '__main__':
    completed = asyncio.run(run())
    print(f'PASS {len(completed)}')
    for index, item in enumerate(completed, start=1):
        print(f'{index:02d}. {item}')
